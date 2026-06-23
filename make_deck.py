"""Generate ParkSight pitch deck (.pptx) — clean, formal, minimalist theme.

Run: python make_deck.py  ->  produces ParkSight_Pitch_Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ----- palette -----
INK        = RGBColor(0x1c, 0x1f, 0x26)
INK_BODY   = RGBColor(0x33, 0x37, 0x40)
INK_MUTED  = RGBColor(0x6b, 0x70, 0x7a)
ACCENT     = RGBColor(0xe2, 0x48, 0x3d)
ACCENT_TINT= RGBColor(0xfb, 0xea, 0xe8)
SUBTLE_BG  = RGBColor(0xf6, 0xf7, 0xf9)
BORDER     = RGBColor(0xe5, 0xe7, 0xeb)
GREEN      = RGBColor(0x2e, 0x7d, 0x32)
AMBER      = RGBColor(0xb7, 0x6e, 0x00)
WHITE      = RGBColor(0xff, 0xff, 0xff)

# Cross-platform-safe typography
HEAD = "Georgia"      # formal serif headings
BODY = "Arial"        # clean sans body
MONO = "Courier New"  # formulas
FONT = BODY

IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deck_assets")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
TOTAL_SLIDES = 15


# ----- helpers -----
def _set_run(r, text, size=14, color=INK_BODY, bold=False, italic=False, font=FONT):
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic


def add_text(slide, x, y, w, h, text, size=14, color=INK_BODY, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size=size, color=color, bold=bold, italic=italic, font=font)
    return tx


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_bullets(slide, x, y, w, h, items, size=14, color=INK_BODY, bullet=True, line_spacing=1.2):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(5)
        if bullet:
            _set_run(p.add_run(), "•  ", size=size, color=ACCENT, bold=True)
        if isinstance(item, str):
            _set_run(p.add_run(), item, size=size, color=color)
        else:
            for chunk in item:
                if isinstance(chunk, str):
                    _set_run(p.add_run(), chunk, size=size, color=color)
                else:
                    text, opts = chunk
                    _set_run(p.add_run(), text,
                             size=opts.get("size", size),
                             color=opts.get("color", color),
                             bold=opts.get("bold", False),
                             italic=opts.get("italic", False))
    return tx


def _soft_shadow(shape, blur=Pt(7), dist=Pt(2.5), alpha=16000, dir_deg=90):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = etree.SubElement(spPr, qn('a:effectLst'))
    shdw = etree.SubElement(eff, qn('a:outerShdw'))
    shdw.set('blurRad', str(int(blur)))
    shdw.set('dist', str(int(dist)))
    shdw.set('dir', str(int(dir_deg * 60000)))
    shdw.set('rotWithShape', '0')
    clr = etree.SubElement(shdw, qn('a:srgbClr'))
    clr.set('val', '1C1F26')
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha)))


def add_header(slide, title, kicker=None):
    add_text(slide, Inches(0.55), Inches(0.42), Inches(12.3), Inches(0.65),
             title, size=27, bold=True, color=INK, font=HEAD)
    if kicker:
        add_text(slide, Inches(0.55), Inches(0.99), Inches(12.3), Inches(0.32),
                 kicker, size=11.5, bold=True, color=ACCENT, font=BODY)
    add_rect(slide, Inches(0.55), Inches(1.34), Inches(12.30), Inches(0.012), BORDER)
    add_rect(slide, Inches(0.55), Inches(1.325), Inches(0.70), Inches(0.035), ACCENT)


def add_footer(slide, page):
    add_rect(slide, Inches(0), Inches(7.22), SLIDE_W, Inches(0.035), ACCENT)
    add_text(slide, Inches(0.55), Inches(7.30), Inches(8), Inches(0.20),
             "ParkSight  ·  Flipgrid.STAR (IIT Kanpur)", size=9, color=INK_MUTED)
    add_text(slide, Inches(11.5), Inches(7.30), Inches(1.3), Inches(0.20),
             f"{page} / {TOTAL_SLIDES}", size=9, color=INK_MUTED, align=PP_ALIGN.RIGHT)


def new_slide(title=None, kicker=None, page=None, footer=True):
    s = prs.slides.add_slide(BLANK)
    if title is not None:
        add_header(s, title, kicker)
    if footer and page is not None:
        add_footer(s, page)
    return s


def add_card(slide, x, y, w, h, heading, body, heading_size=14, body_size=12,
             accent=ACCENT, body_color=INK_BODY):
    box = add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    _soft_shadow(box, alpha=11000)
    add_rect(slide, x, y, Inches(0.055), h, fill=accent)
    add_text(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.34), Inches(0.42),
             heading, size=heading_size, bold=True, color=INK, font=HEAD)
    add_text(slide, x + Inches(0.22), y + Inches(0.62), w - Inches(0.34), h - Inches(0.74),
             body, size=body_size, color=body_color)


def add_stat(slide, x, y, w, h, value, label, color=ACCENT):
    box = add_rect(slide, x, y, w, h, fill=WHITE, line=BORDER)
    _soft_shadow(box, alpha=11000)
    add_text(slide, x, y + Inches(0.20), w, Inches(0.7),
             value, size=27, bold=True, color=color, align=PP_ALIGN.CENTER, font=HEAD)
    add_text(slide, x, y + Inches(0.92), w, Inches(0.5),
             label, size=10.5, color=INK_MUTED, align=PP_ALIGN.CENTER)


def add_image_fit(slide, path, bx, by, bw, bh, border=True, shadow=True, crop=None):
    """Place image inside box, preserving aspect, centered. crop=(L,T,R,B) fractions."""
    pic = slide.shapes.add_picture(path, bx, by)
    nat_w, nat_h = pic.width, pic.height
    if crop:
        cl, ct, cr, cb = crop
        pic.crop_left, pic.crop_top, pic.crop_right, pic.crop_bottom = cl, ct, cr, cb
        eff_w = nat_w * (1 - cl - cr)
        eff_h = nat_h * (1 - ct - cb)
    else:
        eff_w, eff_h = nat_w, nat_h
    scale = min(bw / eff_w, bh / eff_h)
    new_w, new_h = int(eff_w * scale), int(eff_h * scale)
    pic.width, pic.height = new_w, new_h
    pic.left = int(bx + (bw - new_w) / 2)
    pic.top = int(by + (bh - new_h) / 2)
    if border:
        pic.line.color.rgb = BORDER
        pic.line.width = Pt(1.0)
    if shadow:
        _soft_shadow(pic, blur=Pt(9), dist=Pt(3), alpha=20000)
    return pic


def caption(slide, x, y, w, text):
    add_text(slide, x, y, w, Inches(0.26), text, size=9.5, italic=True,
             color=INK_MUTED, align=PP_ALIGN.CENTER)


IMG = lambda name: os.path.join(IMG_DIR, name)

# ============================================================
# Slide 1 — Title  (decluttered)
# ============================================================
s = new_slide(footer=False)
add_rect(s, Inches(0), Inches(0), Inches(0.5), SLIDE_H, fill=ACCENT)

# Hackathon badge (top)
add_bullets(s, Inches(1.13), Inches(0.95), Inches(11.5), Inches(0.34), [
    [("GRIDLOCK HACKATHON 2.0   ·   ROUND 2       ", {"bold": True, "color": INK_MUTED, "size": 12}),
     ("by ", {"color": INK_MUTED, "size": 12}),
     ("Flipkart", {"bold": True, "color": RGBColor(0x28, 0x74, 0xF0), "size": 13.5})],
], bullet=False)

add_text(s, Inches(1.1), Inches(2.15), Inches(11.5), Inches(1.05),
         "ParkSight", size=66, bold=True, color=INK, font=HEAD)
add_text(s, Inches(1.13), Inches(3.15), Inches(11.5), Inches(0.55),
         "Parking-Induced Congestion Intelligence", size=23, color=INK_MUTED, font=HEAD)
add_rect(s, Inches(1.15), Inches(3.88), Inches(0.85), Inches(0.04), ACCENT)
add_text(s, Inches(1.13), Inches(4.05), Inches(11.5), Inches(0.40),
         "Detect   →   Quantify   →   Predict   →   Act", size=15, bold=True, color=ACCENT)

# Team line (airy, no heavy box)
add_text(s, Inches(1.13), Inches(5.55), Inches(11.5), Inches(0.34),
         "Team Flipgrid.STAR   ·   IIT Kanpur", size=15, bold=True, color=INK, font=HEAD)
add_text(s, Inches(1.13), Inches(5.95), Inches(11.5), Inches(0.30),
         "Aaditya Rathi   ·   Ankit Kumar   ·   Priyanshi Agarwal   ·   Shivesh Shukla",
         size=12, color=INK_BODY)
add_rect(s, Inches(1.15), Inches(6.45), Inches(6.6), Inches(0.012), BORDER)
add_text(s, Inches(1.13), Inches(6.58), Inches(11.5), Inches(0.28),
         "Problem statement: Poor Visibility on Parking-Induced Congestion",
         size=10.5, italic=True, color=INK_MUTED)


# ============================================================
# Slide 2 — The Problem
# ============================================================
s = new_slide(title="The Problem", kicker="POOR VISIBILITY ON PARKING-INDUCED CONGESTION", page=2)
add_text(s, Inches(0.55), Inches(1.66), Inches(12.3), Inches(0.6),
         "On-street illegal and spillover parking near commercial areas, metro stations, and events "
         "chokes carriageways and intersections.", size=16, color=INK_BODY, italic=True)
y, h, w, gap = Inches(2.55), Inches(1.65), Inches(3.95), Inches(0.22)
add_card(s, Inches(0.55), y, w, h, "Reactive, not proactive",
         "Enforcement is patrol-based — officers respond after the fact, with no forward plan.",
         heading_size=14, body_size=12)
add_card(s, Inches(0.55)+w+gap, y, w, h, "No visibility on impact",
         "Citation data exists, but there is no map of parking violations vs. their congestion impact.",
         heading_size=14, body_size=12)
add_card(s, Inches(0.55)+2*(w+gap), y, w, h, "Cannot prioritise",
         "Without a ranked, evidence-based target list, scarce enforcement spreads thin across the city.",
         heading_size=14, body_size=12)
add_rect(s, Inches(0.55), Inches(4.65), Inches(12.30), Inches(1.85), fill=ACCENT_TINT, line=ACCENT)
add_text(s, Inches(0.85), Inches(4.82), Inches(11.8), Inches(0.32),
         "PROBLEM STATEMENT DIRECTION", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(5.20), Inches(11.8), Inches(1.20),
         "How can AI-driven parking intelligence detect illegal-parking hotspots and quantify their "
         "impact on traffic flow to enable targeted enforcement?",
         size=19, bold=True, color=INK, font=HEAD)


# ============================================================
# Slide 3 — The Dataset
# ============================================================
s = new_slide(title="The Dataset", kicker="WHAT WE WORKED WITH", page=3)
add_text(s, Inches(0.55), Inches(1.66), Inches(12.3), Inches(0.4),
         "On-street parking-enforcement records from Bengaluru Traffic Police (ASTraM), "
         "November 2023 – April 2024.", size=14, color=INK_BODY)
sy, sh, sw, sx0, sgap = Inches(2.15), Inches(1.30), Inches(2.92), Inches(0.55), Inches(0.18)
add_stat(s, sx0, sy, sw, sh, "292,768", "raw citation records")
add_stat(s, sx0+sw+sgap, sy, sw, sh, "~6 months", "of enforcement activity")
add_stat(s, sx0+2*(sw+sgap), sy, sw, sh, "8 classes", "of vehicle (scooter → HGV)")
add_stat(s, sx0+3*(sw+sgap), sy, sw, sh, "city-wide", "Bengaluru coverage")
add_rect(s, Inches(0.55), Inches(3.75), Inches(12.3), Inches(2.80), fill=SUBTLE_BG, line=BORDER)
add_text(s, Inches(0.80), Inches(3.90), Inches(11.8), Inches(0.3),
         "Each record carries:", size=12, bold=True, color=INK_MUTED)
add_bullets(s, Inches(0.80), Inches(4.25), Inches(11.8), Inches(2.25), [
    [("Geocoded location ", {"bold": True}), ("— latitude / longitude of the violation.", {})],
    [("Vehicle class ", {"bold": True}), ("— scooter, car, autorickshaw, bus, lorry, and more.", {})],
    [("Violation type ", {"bold": True}), ("— main-road parking, near a crossing, double parking, etc.", {})],
    [("Timestamp ", {"bold": True}), ("— creation & modification times, converted UTC → IST before analysis.", {})],
    [("Named junction ", {"bold": True}), ("— links violations to recognised traffic-police junctions.", {})],
    [("Validation status ", {"bold": True}), ("— approved / rejected / unvalidated by reviewing officers.", {})],
], size=13)


# ============================================================
# Slide 4 — Data Engineering
# ============================================================
s = new_slide(title="Data Engineering", kicker="FROM RAW RECORDS TO A TRUSTWORTHY PANEL", page=4)
add_text(s, Inches(0.55), Inches(1.66), Inches(12.3), Inches(0.4),
         "We did not analyse the raw file as-is. Three explicit, defensible cleaning steps.",
         size=13, italic=True, color=INK_MUTED)

cy, ch = Inches(2.12), Inches(1.40)
add_rect(s, Inches(0.55), cy, Inches(12.3), ch, fill=SUBTLE_BG, line=BORDER)
add_rect(s, Inches(0.55), cy, Inches(0.055), ch, fill=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.12), Inches(11.8), Inches(0.32),
         "1.  TIMEZONE CORRECTION  (UTC → IST)", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.50), Inches(11.8), Inches(0.85),
         "All timestamps were UTC. We convert to IST before deriving hour-of-day / day-of-week — "
         "otherwise an 11 AM Bengaluru peak would land at 5:30 AM and every temporal insight would be wrong.",
         size=12, color=INK_BODY)

cy = Inches(3.72); ch = Inches(1.80)
add_rect(s, Inches(0.55), cy, Inches(12.3), ch, fill=SUBTLE_BG, line=BORDER)
add_rect(s, Inches(0.55), cy, Inches(0.055), ch, fill=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.12), Inches(11.8), Inches(0.32),
         "2.  DEDUPLICATION OF BURST CAPTURES  (the “duplicate rows” question)", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.50), Inches(11.8), Inches(1.25),
         "~1.8% of records were duplicate captures of one physical event — distinct ticket IDs sharing the same "
         "vehicle, exact location and timestamp-to-the-second, emitted in ~7-second bursts, with at most one ever "
         "validated-approved. We collapsed each burst to a single event (keeping the approved record), so one parked "
         "vehicle is not counted 2–15× in the impact surface.",
         size=11.5, color=INK_BODY)

cy = Inches(5.72); ch = Inches(1.32)
add_rect(s, Inches(0.55), cy, Inches(12.3), ch, fill=SUBTLE_BG, line=BORDER)
add_rect(s, Inches(0.55), cy, Inches(0.055), ch, fill=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.12), Inches(11.8), Inches(0.32),
         "3.  TIMESTAMP INTEGRITY CHECK", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.80), cy+Inches(0.50), Inches(11.8), Inches(0.75),
         "The citation second-field is constant (a logging artifact), but the modification timestamp carries genuine "
         "sub-minute precision and shows the same morning-concentrated profile — so hour & day-of-week are reliable.",
         size=11.5, color=INK_BODY)


# ============================================================
# Slide 5 — Approach (flow, no number badges)
# ============================================================
s = new_slide(title="Our Approach", kicker="A FOUR-STEP DECISION-SUPPORT PIPELINE", page=5)
add_text(s, Inches(0.55), Inches(1.66), Inches(12.3), Inches(0.4),
         "Turning raw citations into a deployable, forward-looking enforcement plan.",
         size=14, italic=True, color=INK_MUTED)

step_y, step_h, step_w, step_gap = Inches(2.45), Inches(2.70), Inches(2.85), Inches(0.30)
labels = ["DETECT", "QUANTIFY", "PREDICT", "ACT"]
headers = ["Where are real hotspots?", "How bad for traffic flow?",
           "What's coming next week?", "Where & when to deploy?"]
bodies = [
    "Bin citations onto ~65 m H3 hexagons; the Getis-Ord Gi* statistic flags clusters that are statistically significant (Critical / High).",
    "Match each hotspot to its OpenStreetMap road class & lane count; estimate the share of carriageway capacity a stationary vehicle removes.",
    "A gradient-boosted model scores every hotspot × hour slot for the week ahead → per-zone High / Medium / Low risk + predicted peak window.",
    "Output a ranked, downloadable patrol schedule with the forecast attached, plus an interactive day × hour targeting view.",
]
for i in range(4):
    x = Inches(0.55) + i * (step_w + step_gap)
    box = add_rect(s, x, step_y, step_w, step_h, fill=WHITE, line=BORDER)
    _soft_shadow(box, alpha=11000)
    add_rect(s, x, step_y, step_w, Inches(0.10), fill=ACCENT)
    add_text(s, x + Inches(0.22), step_y + Inches(0.28), step_w - Inches(0.40), Inches(0.34),
             labels[i], size=13, bold=True, color=ACCENT, font=HEAD)
    add_text(s, x + Inches(0.22), step_y + Inches(0.72), step_w - Inches(0.40), Inches(0.50),
             headers[i], size=13, bold=True, color=INK)
    add_text(s, x + Inches(0.22), step_y + Inches(1.22), step_w - Inches(0.40), Inches(1.40),
             bodies[i], size=10.5, color=INK_BODY)
    if i < 3:
        ax = x + step_w + Inches(0.03)
        ay = step_y + step_h/2 - Inches(0.11)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.24), Inches(0.22))
        arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background(); arr.shadow.inherit = False

add_rect(s, Inches(0.55), Inches(5.55), Inches(12.30), Inches(1.15), fill=ACCENT_TINT, line=ACCENT)
add_text(s, Inches(0.85), Inches(5.69), Inches(11.7), Inches(0.36),
         "THE OUTCOME", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(6.04), Inches(11.7), Inches(0.60),
         "Reactive patrolling becomes proactive, scheduled deployment — with an ML risk forecast on every priority zone.",
         size=15, bold=True, color=INK, font=HEAD)


# ============================================================
# Slide 6 — Detect  [spatial_map]
# ============================================================
s = new_slide(title="Detect — Real Hotspots", kicker="GETIS-ORD Gi*  ·  STATISTICAL SIGNIFICANCE, NOT JUST VOLUME", page=6)
lx, lw = Inches(0.55), Inches(4.55)
add_bullets(s, lx, Inches(1.72), lw, Inches(3.1), [
    [("65 m H3 grid.  ", {"bold": True}), ("Citations binned to an equal-area hex grid (6,805 cells).", {})],
    [("Getis-Ord Gi*.  ", {"bold": True}), ("A z-score + p-value per cell — flags genuine clustering, not just where officers patrolled.", {})],
    [("Significance tiers.  ", {"bold": True}), ("Critical at 99% confidence; High at 95%.", {})],
], size=13, line_spacing=1.3)
add_stat(s, lx, Inches(5.00), Inches(2.18), Inches(1.30), "412", "Critical zones (99%)")
add_stat(s, lx + Inches(2.37), Inches(5.00), Inches(2.18), Inches(1.30), "93", "High zones (95%)", color=AMBER)
add_image_fit(s, IMG("spatial_map.png"), Inches(5.35), Inches(1.62), Inches(7.45), Inches(5.28),
              crop=(0.02, 0.26, 0.02, 0.20))
caption(s, Inches(5.35), Inches(6.92), Inches(7.45),
        "Live dashboard — Critical (red) & High (amber) ~65 m hotspot cells across Bengaluru.")


# ============================================================
# Slide 7 — Quantify  [methodology_tables]
# ============================================================
s = new_slide(title="Quantify — Traffic-Flow Impact", kicker="ANCHORED TO REAL ROAD GEOMETRY, NOT A HEURISTIC", page=7)
add_rect(s, Inches(0.55), Inches(1.62), Inches(12.30), Inches(1.14), fill=SUBTLE_BG, line=BORDER)
add_rect(s, Inches(0.55), Inches(1.62), Inches(0.055), Inches(1.14), fill=ACCENT)
add_text(s, Inches(0.80), Inches(1.73), Inches(11.8), Inches(0.30),
         "IMPACT SCORE PER CITATION", size=10.5, bold=True, color=ACCENT)
add_text(s, Inches(0.80), Inches(2.04), Inches(11.8), Inches(0.42),
         "impact  =  vehicle footprint  ×  violation severity  ×  time-of-day demand",
         size=16, bold=True, color=INK, font=MONO)
add_text(s, Inches(0.80), Inches(2.44), Inches(11.8), Inches(0.28),
         "lorry > scooter   ·   double-park on a main road > side lane   ·   rush hour ×1.3",
         size=10.5, italic=True, color=INK_MUTED)

lx, lw = Inches(0.55), Inches(4.55)
add_bullets(s, lx, Inches(3.05), lw, Inches(2.4), [
    [("Match to OSM.  ", {"bold": True}), ("Each hotspot is matched to its nearest road for real road class & lane count.", {})],
    [("Capacity loss %.  ", {"bold": True}), ("blocked ÷ total lanes, scaled by vehicle size, amplified at junctions.", {})],
    [("Why it works.  ", {"bold": True}), ("A single-lane blockage hurts far more than the same car on an arterial.", {})],
], size=12.5, line_spacing=1.25)
add_stat(s, lx, Inches(5.55), Inches(2.18), Inches(1.30), "98%", "max capacity removed")
add_stat(s, lx + Inches(2.37), Inches(5.55), Inches(2.18), Inches(1.30), "~48%", "avg in Critical zones", color=AMBER)
add_image_fit(s, IMG("methodology_tables.png"), Inches(5.35), Inches(3.05), Inches(7.45), Inches(3.55))
caption(s, Inches(5.35), Inches(6.64), Inches(7.45),
        "Live dashboard — vehicle-footprint & violation-severity weighting tables.")


# ============================================================
# Slide 8 — Temporal Demand  [temporal_heatmap]  (NEW)
# ============================================================
s = new_slide(title="Temporal Demand — When Pressure Recurs", kicker="STABLE WEEKLY PATTERNS MAKE SCHEDULING POSSIBLE", page=8)
lx, lw = Inches(0.55), Inches(4.55)
add_bullets(s, lx, Inches(1.72), lw, Inches(3.1), [
    [("Recurring, not random.  ", {"bold": True}), ("Congestion pressure repeats on a stable weekly rhythm.", {})],
    [("Mornings dominate.  ", {"bold": True}), ("Impact concentrates in the 08:00–11:00 window.", {})],
    [("This is the signal.  ", {"bold": True}), ("Because the pattern recurs, enforcement can be scheduled in advance — not just dispatched reactively.", {})],
], size=13, line_spacing=1.3)
add_stat(s, lx, Inches(5.05), Inches(4.55), Inches(1.25), "Sunday 10:00", "single busiest city-wide slot")
add_image_fit(s, IMG("temporal_heatmap.png"), Inches(5.35), Inches(1.70), Inches(7.45), Inches(5.10),
              crop=(0.03, 0.49, 0.045, 0.0))
caption(s, Inches(5.35), Inches(6.86), Inches(7.45),
        "Live dashboard — weekly heatmap with hourly & weekday/weekend impact curves.")


# ============================================================
# Slide 9 — Predict  [enforcement_table]
# ============================================================
s = new_slide(title="Predict — Where the AI Lives", kicker="GRADIENT-BOOSTED CLASSIFIER  ·  84,840 PREDICTIONS PER WEEK", page=9)
lx, lw = Inches(0.55), Inches(4.55)
add_bullets(s, lx, Inches(1.72), lw, Inches(3.4), [
    [("Model.  ", {"bold": True}), ("HistGradientBoosting classifier (scikit-learn).", {})],
    [("Predicts.  ", {"bold": True}), ("probability a hotspot × day × hour slot is a top-quartile congestion event (9 features: time, climatology, Gi*, road geometry).", {})],
    [("Next week.  ", {"bold": True}), ("505 hotspots × 7 days × 24 hours = 84,840 scored slots → risk tier + predicted peak window.", {})],
    [("Honest framing.  ", {"bold": True}), ("a ranker, not a magic forecaster; validated on a held-out future period (no leakage).", {})],
], size=12.5, line_spacing=1.22)
add_stat(s, lx, Inches(5.35), Inches(2.18), Inches(1.35), "0.63", "model ROC-AUC")
add_stat(s, lx + Inches(2.37), Inches(5.35), Inches(2.18), Inches(1.35), "+7%", "lift vs naive baseline", color=GREEN)
add_image_fit(s, IMG("enforcement_table.png"), Inches(5.35), Inches(1.62), Inches(7.45), Inches(5.28),
              crop=(0.01, 0.02, 0.01, 0.02))
caption(s, Inches(5.35), Inches(6.92), Inches(7.45),
        "Live dashboard — ranked schedule with per-zone forecast risk, peak probability & window.")


# ============================================================
# Slide 10 — Act  [timewindow_map]
# ============================================================
s = new_slide(title="Act — A Deployable Patrol Plan", kicker="THE END PRODUCT  ·  WHERE & WHEN OFFICERS SHOULD GO", page=10)
lx, lw = Inches(0.55), Inches(4.55)
cy, ch, cgap = Inches(1.72), Inches(1.58), Inches(0.16)
add_card(s, lx, cy, lw, ch, "Ranked worklist",
         "Top 30 hotspots with violation, capacity loss & recurring peak — downloads as CSV into ASTraM's roster.",
         heading_size=13, body_size=11)
add_card(s, lx, cy+ch+cgap, lw, ch, "Next-week risk forecast",
         "Per-zone High / Medium / Low tier + the single most likely peak hour. 169 zones flagged High-risk.",
         heading_size=13, body_size=11)
add_card(s, lx, cy+2*(ch+cgap), lw, ch, "Time-window targeting",
         "Pick any day × hour; the map updates live to show active hotspots for that window.",
         heading_size=13, body_size=11)
add_image_fit(s, IMG("timewindow_map.png"), Inches(5.35), Inches(1.62), Inches(7.45), Inches(5.28),
              crop=(0.02, 0.14, 0.02, 0.30))
caption(s, Inches(5.35), Inches(6.92), Inches(7.45),
        "Live dashboard — interactive day × hour targeting (shown: Sunday 10:00).")


# ============================================================
# Slide 11 — Key Findings  (Pareto; unique stats, no image)
# ============================================================
s = new_slide(title="Key Findings", kicker="A SMALL TEAM CAN COVER MOST OF THE PROBLEM", page=11)
add_rect(s, Inches(0.55), Inches(1.85), Inches(12.30), Inches(1.35), fill=ACCENT_TINT, line=ACCENT)
add_text(s, Inches(0.85), Inches(2.02), Inches(11.7), Inches(0.32),
         "THE PARETO INSIGHT", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(2.38), Inches(11.7), Inches(0.70),
         "The top 5% of locations carry ~60% of all congestion impact — the top 1% alone ≈ 30%.",
         size=22, bold=True, color=INK, font=HEAD)

sy, sh, sw, sx0, sgap = Inches(3.65), Inches(1.55), Inches(3.92), Inches(0.55), Inches(0.18)
add_stat(s, sx0, sy, sw, sh, "60%", "of impact from the top 5% of cells")
add_stat(s, sx0+sw+sgap, sy, sw, sh, "30%", "of impact from the top 1% of cells")
add_stat(s, sx0+2*(sw+sgap), sy, sw, sh, "~7%", "of cells contain most of the problem")

add_bullets(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(1.3), [
    [("Geographic concentration.  ", {"bold": True}), ("Worst intersection: Safina Plaza Junction; worst single cell: Kadubeesanahalli Underpass (~95% capacity loss when blocked).", {})],
    [("Operational takeaway.  ", {"bold": True}), ("Focus a small enforcement team on a few hundred zones and you address the bulk of the city's parking-induced congestion.", {})],
], size=12.5, line_spacing=1.2)


# ============================================================
# Slide 12 — Mapping to Problem Statement  [intersections]
# ============================================================
s = new_slide(title="Mapping to the Problem Statement", kicker="EVERY ASK IN THE BRIEF, ANSWERED", page=12)
lx, lw = Inches(0.55), Inches(5.05)
add_text(s, lx, Inches(1.70), lw, Inches(0.32), "The three core asks", size=13, bold=True, color=INK, font=HEAD)
add_bullets(s, lx, Inches(2.06), lw, Inches(1.55), [
    [("Detect ", {"bold": True, "color": GREEN}), ("→ Gi* significance map + Critical / High tiers.", {})],
    [("Quantify ", {"bold": True, "color": GREEN}), ("→ % carriageway capacity removed, from real OSM geometry.", {})],
    [("Target ", {"bold": True, "color": GREEN}), ("→ ranked schedule + ML next-week risk forecast.", {})],
], size=12.5, line_spacing=1.25)
add_text(s, lx, Inches(3.75), lw, Inches(0.32), "The four named contexts", size=13, bold=True, color=INK, font=HEAD)
add_bullets(s, lx, Inches(4.11), lw, Inches(2.6), [
    [("Commercial / markets ", {"bold": True}), ("— 9.0% of impact.", {})],
    [("Institutional ", {"bold": True}), ("— 5.9% (hospitals / schools / courts).", {})],
    [("Transit (metro / bus / rail) ", {"bold": True}), ("— 1.2%.", {})],
    [("Events / hospitality ", {"bold": True}), ("— 0.5%.", {})],
    [("Intersections ", {"bold": True}), ("— ranked named junctions (top: Safina Plaza).", {})],
], size=12, line_spacing=1.2)
add_image_fit(s, IMG("intersections.png"), Inches(5.85), Inches(1.66), Inches(6.95), Inches(5.20),
              crop=(0.01, 0.115, 0.01, 0.02))
caption(s, Inches(5.85), Inches(6.90), Inches(6.95),
        "Live dashboard — junctions, place types & what drives the impact.")


# ============================================================
# Slide 13 — Why We Trust This
# ============================================================
s = new_slide(title="Why We Trust This", kicker="MAPPED TO THE JUDGING DIMENSIONS", page=13)
cy, cw, ch, cgx, cgy = Inches(1.75), Inches(6.05), Inches(2.32), Inches(0.20), Inches(0.20)
add_card(s, Inches(0.55), cy, cw, ch, "Robust",
         "Statistical significance via Gi* p-values; the ML model is validated on a held-out future period and beats a seasonal-naive baseline. We also report an honest negative result.",
         heading_size=15, body_size=12)
add_card(s, Inches(0.55)+cw+cgx, cy, cw, ch, "Innovative",
         "Fuses three methods rarely combined: spatial statistics (Gi*), OSM road-geometry capacity loss, and an ML next-week risk forecast — predictive prioritisation, not after-the-fact analytics.",
         heading_size=15, body_size=12)
add_card(s, Inches(0.55), cy+ch+cgy, cw, ch, "Scalable",
         "City- and map-provider-agnostic; the whole pipeline runs on a standard laptop in under 20 minutes per city, with no city-specific tuning.",
         heading_size=15, body_size=12)
add_card(s, Inches(0.55)+cw+cgx, cy+ch+cgy, cw, ch, "Real-world",
         "Outputs are downloadable CSV schedules at the per-hour, per-day cadence ASTraM already operates on — officers download and drop them into the roster.",
         heading_size=15, body_size=12)


# ============================================================
# Slide 14 — Limitations & Future Work  (merged)
# ============================================================
s = new_slide(title="Limitations & Future Work", kicker="WHAT WE DON'T CLAIM — AND HOW WE'D EXTEND IT", page=14)
colw = Inches(6.05); lxx = Inches(0.55); rxx = Inches(0.55)+colw+Inches(0.20)
add_text(s, lxx, Inches(1.68), colw, Inches(0.34), "Honest limitations", size=14, bold=True, color=ACCENT, font=HEAD)
add_bullets(s, lxx, Inches(2.08), colw, Inches(4.7), [
    [("Enforcement-data proxy.  ", {"bold": True}), ("Records show where officers ticketed — a proxy for, not a census of, illegal parking (patrol selection bias).", {})],
    [("No live traffic signal.  ", {"bold": True}), ("Capacity loss is modelled from static OSM geometry, not measured from observed speeds.", {})],
    [("No dwell time.  ", {"bold": True}), ("The strongest physical determinant of impact is absent; we proxy it with footprint × severity.", {})],
    [("Uncalibrated probabilities.  ", {"bold": True}), ("Forecast scores rank well but aren't calibrated.", {})],
], size=12, line_spacing=1.25)
add_text(s, rxx, Inches(1.68), colw, Inches(0.34), "How we'd extend it", size=14, bold=True, color=GREEN, font=HEAD)
add_bullets(s, rxx, Inches(2.08), colw, Inches(4.7), [
    [("Live speed feeds.  ", {"bold": True}), ("Fuse probe-vehicle / ASTraM speed data to calibrate capacity loss against real flow degradation.", {})],
    [("Patrol-exposure debiasing.  ", {"bold": True}), ("Normalise impact by per-station patrol exposure to remove selection bias.", {})],
    [("Causal evaluation.  ", {"bold": True}), ("Difference-in-differences around deployment changes to prove enforcement worked.", {})],
    [("MapmyIndia road graph.  ", {"bold": True}), ("Swap OSM for higher-resolution Indian lane data — no retraining.", {})],
], size=12, line_spacing=1.25)


# ============================================================
# Slide 15 — Partners, Scalability & Team  (close)
# ============================================================
s = new_slide(title="Partners, Scalability & Team", kicker="BUILT TO PLUG IN, BUILT TO TRAVEL", page=15)
py, ph, pw, pgap = Inches(1.70), Inches(2.05), Inches(6.05), Inches(0.20)
add_card(s, Inches(0.55), py, pw, ph, "ASTraM (Bengaluru Traffic Police)",
         "A decision-support layer for the existing operational stack. Field-ready CSV schedules at ASTraM's per-hour, per-day cadence drop straight into the patrol roster — no new system, no new infrastructure.",
         heading_size=14, body_size=11.5)
add_card(s, Inches(0.55)+pw+pgap, py, pw, ph, "MapmyIndia",
         "A drop-in upgrade for the road-geometry layer. We use OpenStreetMap today; MapmyIndia's proprietary lane graph slots in with no model retraining — a code-level swap of the geometry source.",
         heading_size=14, body_size=11.5)
add_rect(s, Inches(0.55), Inches(3.95), Inches(12.30), Inches(0.78), fill=SUBTLE_BG, line=BORDER)
add_rect(s, Inches(0.55), Inches(3.95), Inches(0.055), Inches(0.78), fill=ACCENT)
add_text(s, Inches(0.80), Inches(4.07), Inches(11.8), Inches(0.55),
         "Scalable — city- & map-provider-agnostic; any city with geocoded citations + OpenStreetMap runs it unchanged.",
         size=13, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.55), Inches(5.05), Inches(12.30), Inches(1.45), fill=WHITE, line=BORDER)
add_text(s, Inches(0.80), Inches(5.18), Inches(11.8), Inches(0.32),
         "TEAM FLIPGRID.STAR  ·  IIT KANPUR", size=11, bold=True, color=ACCENT)
add_text(s, Inches(0.80), Inches(5.52), Inches(11.8), Inches(0.42),
         "Aaditya Rathi   ·   Ankit Kumar   ·   Priyanshi Agarwal   ·   Shivesh Shukla",
         size=15, bold=True, color=INK, font=HEAD)
add_text(s, Inches(0.80), Inches(6.05), Inches(11.8), Inches(0.32),
         "Live dashboard: parksight-flipgridstar.streamlit.app    ·    Code: github.com/arathii23/parksight",
         size=10.5, color=INK_BODY)


# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ParkSight_Pitch_Deck.pptx")
prs.save(out)
print(f"wrote {out}  ({os.path.getsize(out)/1024:.1f} KB, slides={len(prs.slides)})")
